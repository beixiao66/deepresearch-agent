from pathlib import Path

from app.services.document_parser import DocumentParser


def test_parse_docx_document(tmp_path) -> None:
    from docx import Document as DocxDocument

    document_path = tmp_path / "notes.docx"

    document = DocxDocument()
    document.add_paragraph("第一条段落")
    document.add_paragraph("第二条段落")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "姓名"
    table.cell(0, 1).text = "职位"
    table.cell(1, 0).text = "张三"
    table.cell(1, 1).text = "工程师"
    document.save(str(document_path))

    parsed = DocumentParser().parse(
        str(document_path),
        ".docx",
    )

    assert "第一条段落" in parsed.text
    assert "第二条段落" in parsed.text
    assert "姓名" in parsed.text
    assert "工程师" in parsed.text
    assert parsed.page_count is None


def test_parse_html_document(tmp_path) -> None:
    document_path = tmp_path / "page.html"
    document_path.write_text(
        (
            "<html><head><title>测试</title></head><body>"
            "<script>alert('x')</script>"
            "<nav>导航</nav>"
            "<h1>标题一</h1>"
            "<p>正文内容</p>"
            "</body></html>"
        ),
        encoding="utf-8",
    )

    parsed = DocumentParser().parse(
        str(document_path),
        ".html",
    )

    assert "标题一" in parsed.text
    assert "正文内容" in parsed.text
    assert "alert" not in parsed.text
    assert "导航" not in parsed.text


def test_parse_xlsx_document(tmp_path) -> None:
    from openpyxl import Workbook

    document_path = tmp_path / "data.xlsx"

    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "成绩表"
    worksheet.append(["姓名", "分数"])
    worksheet.append(["张三", 90])
    workbook.save(str(document_path))

    parsed = DocumentParser().parse(
        str(document_path),
        ".xlsx",
    )

    assert "成绩表" in parsed.text
    assert "姓名" in parsed.text
    assert "分数" in parsed.text
    assert "张三" in parsed.text
    assert "90" in parsed.text


def test_parse_pptx_document(tmp_path) -> None:
    from pptx import Presentation

    document_path = tmp_path / "slides.pptx"

    presentation = Presentation()
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[0]
    )
    slide.shapes.title.text = "第一页标题"
    presentation.save(str(document_path))

    parsed = DocumentParser().parse(
        str(document_path),
        ".pptx",
    )

    assert "第一页" in parsed.text
    assert "第一页标题" in parsed.text


def test_parse_csv_document(tmp_path) -> None:
    document_path = tmp_path / "data.csv"
    document_path.write_text(
        "姓名,职位\n张三,工程师\n李四,产品经理\n",
        encoding="utf-8",
    )

    parsed = DocumentParser().parse(
        str(document_path),
        ".csv",
    )

    assert "姓名" in parsed.text
    assert "职位" in parsed.text
    assert "张三" in parsed.text
    assert "工程师" in parsed.text
    assert "李四" in parsed.text
