from dataclasses import dataclass

import logging
from pathlib import Path

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from app.services.file_storage import ALLOWED_DOCUMENT_EXTENSIONS


@dataclass(frozen=True)
class ParsedDocument:
    text: str
    page_count: int | None = None


logger = logging.getLogger(__name__)


class DocumentParser:
    def parse(
            self,
            storage_path: str,
            file_extension: str,
    ) -> ParsedDocument:
        file_path = Path(storage_path)

        if file_extension == ".txt":
            return self._parse_text(file_path)
        if file_extension == ".md":
            return self._parse_text(file_path)
        if file_extension == ".pdf":
            return self._parse_pdf(file_path)
        if file_extension == ".docx":
            return self._parse_docx(file_path)
        if file_extension in {".html", ".htm"}:
            return self._parse_html(file_path)
        if file_extension == ".xlsx":
            return self._parse_xlsx(file_path)
        if file_extension == ".pptx":
            return self._parse_pptx(file_path)
        if file_extension == ".csv":
            return self._parse_csv(file_path)

        raise ValueError(
            f"Unsupported document extension: {file_extension}"
        )

    @staticmethod
    def _parse_text(file_path: Path) -> ParsedDocument:
        text = file_path.read_text(encoding="utf-8")
        return ParsedDocument(text=text)

    @staticmethod
    def _parse_pdf(file_path: Path) -> ParsedDocument:
        reader = PdfReader(str(file_path))

        pages = []
        for page in reader.pages:
            pages.append(page.extract_text() or "")

        text = "\n".join(pages).strip()

        return ParsedDocument(
            text=text,
            page_count=len(reader.pages),
        )

    @staticmethod
    def _parse_docx(file_path: Path) -> ParsedDocument:
        """解析 Word 文档：段落 + 表格，按原文顺序提取。"""
        document = DocxDocument(str(file_path))

        parts: list[str] = []

        # python-docx 的 body 元素按文档顺序混排段落与表格
        from docx.document import Document as DocxDocumentType
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        body = document.element.body

        for child in body.iterchildren():
            if child.tag == qn("w:p"):
                paragraph = Paragraph(child, document)
                text = paragraph.text.strip()
                if text:
                    parts.append(text)
            elif child.tag == qn("w:tbl"):
                table = Table(child, document)
                for row in table.rows:
                    cells = [
                        cell.text.strip()
                        for cell in row.cells
                    ]
                    row_text = " | ".join(
                        cell for cell in cells if cell
                    )
                    if row_text:
                        parts.append(row_text)

        return ParsedDocument(text="\n".join(parts).strip())

    @staticmethod
    def _parse_html(file_path: Path) -> ParsedDocument:
        """解析 HTML：去除脚本/样式/导航后提取正文文本。"""
        content = file_path.read_text(
            encoding="utf-8",
            errors="replace",
        )
        soup = BeautifulSoup(content, "lxml")

        for element in soup(
            [
                "script",
                "style",
                "noscript",
                "nav",
                "header",
                "footer",
                "aside",
            ]
        ):
            element.decompose()

        text = soup.get_text(separator="\n", strip=True)

        return ParsedDocument(text=text)

    @staticmethod
    def _parse_xlsx(file_path: Path) -> ParsedDocument:
        """解析 Excel：工作表 → 行 → 单元格，保留单元格文本。"""
        workbook = load_workbook(
            str(file_path),
            read_only=True,
            data_only=True,
        )

        parts: list[str] = []

        for worksheet in workbook.worksheets:
            parts.append(f"【工作表：{worksheet.title}】")

            for row in worksheet.iter_rows():
                cells = [
                    str(cell.value).strip()
                    for cell in row
                    if cell.value is not None
                ]
                if cells:
                    parts.append(" | ".join(cells))

        workbook.close()

        return ParsedDocument(text="\n".join(parts).strip())

    @staticmethod
    def _parse_pptx(file_path: Path) -> ParsedDocument:
        """解析 PPT：按页提取所有文本框内容。"""
        presentation = Presentation(str(file_path))

        parts: list[str] = []

        for slide_index, slide in enumerate(
                presentation.slides,
                start=1,
        ):
            slide_texts: list[str] = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                text = "\n".join(
                    paragraph.text.strip()
                    for paragraph in shape.text_frame.paragraphs
                    if paragraph.text.strip()
                )
                if text:
                    slide_texts.append(text)

            if slide_texts:
                parts.append(
                    f"【第 {slide_index} 页】\n"
                    + "\n".join(slide_texts)
                )

        return ParsedDocument(text="\n".join(parts).strip())

    @staticmethod
    def _parse_csv(file_path: Path) -> ParsedDocument:
        """解析 CSV：逐行读取，行内单元格用 | 连接。"""
        import csv

        parts: list[str] = []

        with file_path.open(
            "r",
            encoding="utf-8",
            newline="",
            errors="replace",
        ) as file:
            reader = csv.reader(file)

            for row in reader:
                row_text = " | ".join(
                    cell.strip()
                    for cell in row
                    if cell.strip()
                )
                if row_text:
                    parts.append(row_text)

        return ParsedDocument(text="\n".join(parts).strip())